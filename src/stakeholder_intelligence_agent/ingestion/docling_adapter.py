"""Docling-primary six-format extraction with fixture-proven narrow supplements."""

from __future__ import annotations

import importlib.util
import io
import json
import math
import os
import shutil
import textwrap
from collections import defaultdict, deque
from dataclasses import replace
from pathlib import Path
from typing import TYPE_CHECKING, Any, Final

import pypdfium2 as pdfium
from docling.datamodel.base_models import ConversionStatus, InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions, RapidOcrOptions
from docling.document_converter import (
    DocumentConverter,
    ImageFormatOption,
    PdfFormatOption,
)
from docling.models.stages.layout.layout_model import LayoutModel
from docling.models.stages.table_structure.table_structure_model import TableStructureModel
from docling_core.types.doc.items.picture.picture import PictureItem
from docling_core.types.doc.items.table.table import TableItem
from docling_core.types.doc.items.text import TextItem
from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_to_tuple
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from stakeholder_intelligence_agent.contracts.source import (
    BoundingBox,
    DocxRenderedPageLocation,
    ImageRegionLocation,
    PdfPageLocation,
    PptxSlideLocation,
    XlsxRangeLocation,
)
from stakeholder_intelligence_agent.errors import (
    ExtractionFailedError,
    MandatoryContentMissingError,
)
from stakeholder_intelligence_agent.ingestion.types import (
    ArtifactDraft,
    ElementDraft,
    ExtractionBundle,
    ValidatedUpload,
)

if TYPE_CHECKING:
    from collections.abc import Sequence

    from docling_core.types.doc.common.reference import ProvenanceItem
    from docling_core.types.doc.document import DoclingDocument

    from stakeholder_intelligence_agent.config import Settings
    from stakeholder_intelligence_agent.ingestion.types import ArtifactKind

_CHART_COLORS: Final[tuple[str, ...]] = ("#0B7285", "#F59F00", "#5F3DC4", "#C92A2A")
_TEXT_RENDER_BOTTOM = 2_120


class DoclingExtractor:
    """Run Docling for every format and normalize its source lineage."""

    def __init__(self, settings: Settings) -> None:
        self._settings = settings
        self._converter: DocumentConverter | None = None

    def extract(self, source_path: Path, upload: ValidatedUpload) -> ExtractionBundle:
        """Convert one preserved original and return typed deterministic drafts."""
        result = self._get_converter().convert(source_path)
        if result.status != ConversionStatus.SUCCESS:
            raise ExtractionFailedError
        document = result.document
        if upload.document_type == "pdf":
            bundle = self._extract_pdf(document, source_path, upload)
        elif upload.document_type == "docx":
            bundle = self._extract_docx(document, upload)
        elif upload.document_type == "pptx":
            bundle = self._extract_pptx(document, source_path, upload)
        elif upload.document_type == "xlsx":
            bundle = self._extract_xlsx(document, source_path, upload)
        else:
            bundle = self._extract_image(document, upload)
        if not bundle.elements or not any(
            element.original_content and element.original_content.strip()
            for element in bundle.elements
        ):
            raise MandatoryContentMissingError
        return bundle

    def _get_converter(self) -> DocumentConverter:
        if self._converter is not None:
            return self._converter
        cache_root = self._settings.model_cache_root.resolve()
        cache_root.mkdir(parents=True, exist_ok=True)
        os.environ["HF_HOME"] = str(cache_root / "huggingface")
        os.environ["XDG_CACHE_HOME"] = str(cache_root)
        artifacts_path = self._prepare_model_artifacts(cache_root)
        pipeline_options = PdfPipelineOptions(
            do_ocr=True,
            ocr_options=self._rapidocr_options(),
            do_table_structure=True,
            generate_page_images=True,
            generate_picture_images=True,
            images_scale=1.5,
            document_timeout=float(self._settings.docling_timeout_seconds),
            enable_remote_services=False,
            allow_external_plugins=False,
            artifacts_path=artifacts_path,
        )
        self._converter = DocumentConverter(
            allowed_formats=[
                InputFormat.PDF,
                InputFormat.DOCX,
                InputFormat.PPTX,
                InputFormat.XLSX,
                InputFormat.IMAGE,
            ],
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                InputFormat.IMAGE: ImageFormatOption(pipeline_options=pipeline_options),
            },
        )
        return self._converter

    @classmethod
    def _prepare_model_artifacts(cls, cache_root: Path) -> Path:
        artifacts_root = cache_root / "docling-artifacts"
        layout_target = artifacts_root / "docling-project--docling-layout-heron"
        table_target = artifacts_root / "docling-project--docling-models"
        layout_marker = layout_target / "model.safetensors"
        table_marker = table_target / "model_artifacts" / "tableformer" / "fast" / "tm_config.json"
        if not layout_marker.is_file() and not cls._restore_huggingface_snapshot(
            cache_root,
            repository_cache_name="models--docling-project--docling-layout-heron",
            target=layout_target,
        ):
            LayoutModel.download_models(local_dir=layout_target)
        if not table_marker.is_file() and not cls._restore_huggingface_snapshot(
            cache_root,
            repository_cache_name="models--docling-project--docling-models",
            target=table_target,
        ):
            TableStructureModel.download_models(local_dir=table_target)
        if not layout_marker.is_file() or not table_marker.is_file():
            raise ExtractionFailedError
        return artifacts_root

    @staticmethod
    def _rapidocr_options() -> RapidOcrOptions:
        specification = importlib.util.find_spec("rapidocr")
        if specification is None or specification.origin is None:
            raise ExtractionFailedError
        models = Path(specification.origin).resolve().parent / "models"
        font_candidates = (
            Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "arial.ttf",
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf"),
            Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        )
        font_path = next((path for path in font_candidates if path.is_file()), None)
        det_model = models / "PP-OCRv6_det_small.onnx"
        cls_model = models / "ch_ppocr_mobile_v2.0_cls_mobile.onnx"
        rec_model = models / "PP-OCRv6_rec_small.onnx"
        if font_path is None or not all(
            path.is_file() for path in (det_model, cls_model, rec_model)
        ):
            raise ExtractionFailedError
        return RapidOcrOptions(
            lang=["english"],
            force_full_page_ocr=True,
            det_model_path=str(det_model),
            cls_model_path=str(cls_model),
            rec_model_path=str(rec_model),
            font_path=str(font_path),
        )

    @classmethod
    def _restore_huggingface_snapshot(
        cls,
        cache_root: Path,
        *,
        repository_cache_name: str,
        target: Path,
    ) -> bool:
        snapshots = cache_root / "huggingface" / "hub" / repository_cache_name / "snapshots"
        candidates = sorted(
            (path for path in snapshots.glob("*") if path.is_dir()),
            key=lambda path: path.stat().st_mtime_ns,
            reverse=True,
        )
        if not candidates:
            return False
        target.mkdir(parents=True, exist_ok=True)
        shutil.copytree(
            candidates[0],
            target,
            dirs_exist_ok=True,
            copy_function=cls._hardlink_or_copy,
        )
        return True

    @staticmethod
    def _hardlink_or_copy(source: str, target: str) -> str:
        try:
            os.link(source, target)
        except OSError:
            shutil.copy2(source, target)
        return target

    def _extract_pdf(
        self,
        document: DoclingDocument,
        source_path: Path,
        upload: ValidatedUpload,
    ) -> ExtractionBundle:
        native_pages = self._pdf_native_text_pages(source_path)
        elements: list[ElementDraft] = []
        artifacts: list[ArtifactDraft] = []
        scanned_pages = {page for page, has_text in native_pages.items() if not has_text}
        for page_number in sorted(scanned_pages):
            page = document.pages.get(page_number)
            image = None if page is None or page.image is None else page.image.pil_image
            if image is None:
                continue
            artifact_key = f"pdf-page-render-{page_number}"
            artifacts.append(
                ArtifactDraft(
                    key=artifact_key,
                    artifact_kind="page_render",
                    media_type="image/png",
                    suffix=".png",
                    content=self._png_bytes(image),
                )
            )
            elements.append(
                ElementDraft(
                    key=f"pdf-page-image-{page_number}",
                    element_type="image",
                    original_content=None,
                    location=PdfPageLocation(filename=upload.filename, page=page_number),
                    extraction_method="docling_page_render_v1",
                    artifact_key=artifact_key,
                )
            )

        for item, _level in document.iterate_items():
            provenance = self._first_provenance(getattr(item, "prov", ()))
            page_number = 1 if provenance is None else provenance.page_no
            location = PdfPageLocation(
                filename=upload.filename,
                page=page_number,
                bounding_box=self._bbox(provenance, scale=1.0),
            )
            if isinstance(item, TextItem):
                content = item.text.strip()
                if not content:
                    continue
                is_ocr = page_number in scanned_pages
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="ocr_text" if is_ocr else "text",
                        original_content=content,
                        location=location,
                        extraction_method=(
                            "docling_rapidocr_full_page_v1" if is_ocr else "docling_v2"
                        ),
                        parent_key=f"pdf-page-image-{page_number}" if is_ocr else None,
                    )
                )
            elif isinstance(item, TableItem):
                content = item.export_to_markdown(document).strip()
                if content:
                    elements.append(
                        ElementDraft(
                            key=item.self_ref,
                            element_type="table",
                            original_content=content,
                            location=location,
                            extraction_method="docling_v2",
                        )
                    )
            elif isinstance(item, PictureItem):
                image = item.get_image(document)
                if image is None:
                    continue
                artifact_key = f"{item.self_ref}-visual"
                artifacts.append(
                    ArtifactDraft(
                        key=artifact_key,
                        artifact_kind=(
                            "chart_render" if str(item.label) == "chart" else "embedded_image"
                        ),
                        media_type="image/png",
                        suffix=".png",
                        content=self._png_bytes(image),
                    )
                )
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="chart" if str(item.label) == "chart" else "image",
                        original_content=None,
                        location=location,
                        extraction_method="docling_v2",
                        artifact_key=artifact_key,
                    )
                )
        return ExtractionBundle(
            elements=tuple(elements),
            artifacts=tuple(artifacts),
            capability_facts=(
                "docling_primary",
                "rapidocr_english_force_full_page",
                "pdf_page_and_picture_images",
            ),
        )

    def _extract_image(
        self,
        document: DoclingDocument,
        upload: ValidatedUpload,
    ) -> ExtractionBundle:
        location = ImageRegionLocation(
            filename=upload.filename,
            image_index=1,
            region="whole_image",
            bounding_box=BoundingBox(
                x0=0,
                y0=0,
                x1=1,
                y1=1,
                coordinate_space="normalized",
            ),
        )
        elements: list[ElementDraft] = [
            ElementDraft(
                key="whole-image",
                element_type="image",
                original_content=None,
                location=location,
                extraction_method="preserved_original_v1",
                artifact_key="$original",
            )
        ]
        for item, _level in document.iterate_items():
            if isinstance(item, TextItem) and item.text.strip():
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="ocr_text",
                        original_content=item.text.strip(),
                        location=location,
                        extraction_method="docling_rapidocr_full_page_v1",
                        parent_key="whole-image",
                    )
                )
        return ExtractionBundle(
            elements=tuple(elements),
            artifacts=(),
            capability_facts=(
                "docling_primary",
                "rapidocr_english_force_full_page",
                "whole_image_original_preserved",
            ),
        )

    def _extract_docx(
        self,
        document: DoclingDocument,
        upload: ValidatedUpload,
    ) -> ExtractionBundle:
        elements: list[ElementDraft] = []
        artifacts: list[ArtifactDraft] = []
        render_entries: list[tuple[str, str]] = []
        current_section: str | None = None
        paragraph_number = 0
        for item, _level in document.iterate_items():
            paragraph_number += 1
            if isinstance(item, TextItem):
                content = item.text.strip()
                if not content:
                    continue
                if str(item.label) == "section_header":
                    current_section = content[:500]
                render_entries.append((item.self_ref, content))
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="text",
                        original_content=content,
                        location=DocxRenderedPageLocation(
                            filename=upload.filename,
                            rendered_page=1,
                            section=current_section,
                            paragraph=paragraph_number,
                        ),
                        extraction_method="docling_v2",
                    )
                )
            elif isinstance(item, TableItem):
                content = item.export_to_markdown(document).strip()
                if not content:
                    continue
                render_entries.append((item.self_ref, content))
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="table",
                        original_content=content,
                        location=DocxRenderedPageLocation(
                            filename=upload.filename,
                            rendered_page=1,
                            section=current_section,
                            paragraph=paragraph_number,
                        ),
                        extraction_method="docling_v2",
                    )
                )
            elif isinstance(item, PictureItem):
                image = item.get_image(document)
                if image is None:
                    continue
                artifact_key = f"{item.self_ref}-visual"
                artifacts.append(
                    ArtifactDraft(
                        key=artifact_key,
                        artifact_kind="embedded_image",
                        media_type="image/png",
                        suffix=".png",
                        content=self._png_bytes(image),
                    )
                )
                render_entries.append((item.self_ref, f"[Embedded image: {artifact_key}]"))
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="image",
                        original_content=None,
                        location=DocxRenderedPageLocation(
                            filename=upload.filename,
                            rendered_page=1,
                            section=current_section,
                            paragraph=paragraph_number,
                        ),
                        extraction_method="docling_v2",
                        artifact_key=artifact_key,
                    )
                )
        page_by_element, rendered_pages = self._render_text_pages(
            render_entries,
            title=upload.filename,
        )
        elements = [
            replace(
                element,
                location=element.location.model_copy(
                    update={"rendered_page": page_by_element[element.key]}
                ),
            )
            for element in elements
        ]
        artifacts.extend(
            ArtifactDraft(
                key=f"docx-normalized-render-page-{page_number}",
                artifact_kind="normalized_render",
                media_type="image/png",
                suffix=".png",
                content=content,
            )
            for page_number, content in enumerate(rendered_pages, start=1)
        )
        return ExtractionBundle(
            elements=tuple(elements),
            artifacts=tuple(artifacts),
            capability_facts=(
                "docling_primary",
                "docx_normalized_page_render_supplement_v2",
                "docx_embedded_images",
            ),
        )

    def _extract_pptx(
        self,
        document: DoclingDocument,
        source_path: Path,
        upload: ValidatedUpload,
    ) -> ExtractionBundle:
        chart_renders = self._pptx_chart_renders(source_path)
        elements: list[ElementDraft] = []
        artifacts: list[ArtifactDraft] = []
        for item, _level in document.iterate_items():
            provenance = self._first_provenance(getattr(item, "prov", ()))
            if provenance is None:
                continue
            slide_number = provenance.page_no
            location = PptxSlideLocation(
                filename=upload.filename,
                slide=slide_number,
                shape_identifier=item.self_ref,
                bounding_box=self._bbox(provenance, scale=12_700.0),
            )
            if isinstance(item, TextItem):
                if item.text.strip():
                    elements.append(
                        ElementDraft(
                            key=item.self_ref,
                            element_type="text",
                            original_content=item.text.strip(),
                            location=location,
                            extraction_method="docling_v2",
                        )
                    )
            elif isinstance(item, TableItem):
                content = item.export_to_markdown(document).strip()
                if content:
                    elements.append(
                        ElementDraft(
                            key=item.self_ref,
                            element_type="table",
                            original_content=content,
                            location=location,
                            extraction_method="docling_v2",
                        )
                    )
            elif isinstance(item, PictureItem):
                image = item.get_image(document)
                is_chart = image is None and bool(chart_renders[slide_number])
                artifact_kind: ArtifactKind
                if is_chart:
                    shape_name, visual_content = chart_renders[slide_number].popleft()
                    location = location.model_copy(update={"shape_identifier": shape_name})
                    method = "python_pptx_native_chart_render_v1"
                    artifact_kind = "chart_render"
                elif image is not None:
                    visual_content = self._png_bytes(image)
                    method = "docling_v2"
                    artifact_kind = "embedded_image"
                else:
                    continue
                artifact_key = f"{item.self_ref}-visual"
                artifacts.append(
                    ArtifactDraft(
                        key=artifact_key,
                        artifact_kind=artifact_kind,
                        media_type="image/png",
                        suffix=".png",
                        content=visual_content,
                    )
                )
                elements.append(
                    ElementDraft(
                        key=item.self_ref,
                        element_type="chart" if is_chart else "image",
                        original_content=None,
                        location=location,
                        extraction_method=method,
                        artifact_key=artifact_key,
                    )
                )
        return ExtractionBundle(
            elements=tuple(elements),
            artifacts=tuple(artifacts),
            capability_facts=(
                "docling_primary",
                "pptx_slide_provenance",
                "python_pptx_native_chart_render_supplement_v1",
            ),
        )

    def _extract_xlsx(
        self,
        document: DoclingDocument,
        source_path: Path,
        upload: ValidatedUpload,
    ) -> ExtractionBundle:
        workbook = load_workbook(source_path, data_only=False, read_only=False)
        values_workbook = load_workbook(source_path, data_only=True, read_only=False)
        try:
            if sum(len(sheet._cells) for sheet in workbook.worksheets) > (
                self._settings.max_spreadsheet_cells
            ):
                raise MandatoryContentMissingError
            sheet_names = workbook.sheetnames
            visible_pages = {
                index
                for index, sheet in enumerate(workbook.worksheets, start=1)
                if sheet.sheet_state == "visible"
            }
            chart_renders = self._xlsx_chart_renders(workbook, values_workbook)
            manifest, formula_drafts = self._workbook_manifest_and_formulas(
                workbook,
                values_workbook,
                upload.filename,
            )
            manifest_bytes = (json.dumps(manifest, indent=2, sort_keys=True) + "\n").encode()
            artifacts: list[ArtifactDraft] = [
                ArtifactDraft(
                    key="xlsx-workbook-manifest",
                    artifact_kind="workbook_manifest",
                    media_type="application/json",
                    suffix=".json",
                    content=manifest_bytes,
                )
            ]
            elements: list[ElementDraft] = []
            for item, _level in document.iterate_items():
                provenance = self._first_provenance(getattr(item, "prov", ()))
                if provenance is None or provenance.page_no not in visible_pages:
                    continue
                page_number = provenance.page_no
                if page_number > len(sheet_names):
                    continue
                sheet_name = sheet_names[page_number - 1]
                cell_range = self._xlsx_range(provenance)
                if isinstance(item, TableItem):
                    content = item.export_to_markdown(document).strip()
                    if content:
                        elements.append(
                            ElementDraft(
                                key=item.self_ref,
                                element_type="table",
                                original_content=content,
                                location=XlsxRangeLocation(
                                    filename=upload.filename,
                                    sheet=sheet_name,
                                    cell_range=cell_range,
                                ),
                                extraction_method="docling_v2",
                            )
                        )
                elif isinstance(item, TextItem) and item.text.strip():
                    elements.append(
                        ElementDraft(
                            key=item.self_ref,
                            element_type="text",
                            original_content=item.text.strip(),
                            location=XlsxRangeLocation(
                                filename=upload.filename,
                                sheet=sheet_name,
                                cell_range=cell_range,
                            ),
                            extraction_method="docling_v2",
                        )
                    )
                elif isinstance(item, PictureItem):
                    image = item.get_image(document)
                    is_chart = image is None and bool(chart_renders[sheet_name])
                    artifact_kind: ArtifactKind
                    if is_chart:
                        chart_name, visual_content = chart_renders[sheet_name].popleft()
                        artifact_kind = "chart_render"
                        method = "openpyxl_native_chart_render_v1"
                    elif image is not None:
                        chart_name = ""
                        visual_content = self._png_bytes(image)
                        artifact_kind = "embedded_image"
                        method = "docling_v2"
                    else:
                        continue
                    artifact_key = f"{item.self_ref}-visual"
                    artifacts.append(
                        ArtifactDraft(
                            key=artifact_key,
                            artifact_kind=artifact_kind,
                            media_type="image/png",
                            suffix=".png",
                            content=visual_content,
                        )
                    )
                    elements.append(
                        ElementDraft(
                            key=item.self_ref,
                            element_type="chart" if is_chart else "image",
                            original_content=None,
                            location=XlsxRangeLocation(
                                filename=upload.filename,
                                sheet=sheet_name,
                                cell_range=cell_range,
                                chart_identifier=chart_name or None,
                                image_identifier=None if is_chart else item.self_ref,
                            ),
                            extraction_method=method,
                            artifact_key=artifact_key,
                        )
                    )
            elements.extend(formula_drafts)
            return ExtractionBundle(
                elements=tuple(elements),
                artifacts=tuple(artifacts),
                capability_facts=(
                    "docling_primary",
                    "openpyxl_sheet_formula_visibility_manifest_v1",
                    "openpyxl_native_chart_render_v1",
                    "hidden_sheets_excluded_from_retrieval",
                    "stored_nonempty_cell_bound_enforced",
                ),
            )
        finally:
            workbook.close()
            values_workbook.close()

    @staticmethod
    def _pdf_native_text_pages(source_path: Path) -> dict[int, bool]:
        try:
            document = pdfium.PdfDocument(source_path)
            try:
                return {
                    index + 1: bool(document[index].get_textpage().get_text_range().strip())
                    for index in range(len(document))
                }
            finally:
                document.close()
        except Exception as error:
            raise ExtractionFailedError from error

    @staticmethod
    def _first_provenance(items: Sequence[ProvenanceItem]) -> ProvenanceItem | None:
        return items[0] if items else None

    @staticmethod
    def _bbox(provenance: ProvenanceItem | None, *, scale: float) -> BoundingBox | None:
        if provenance is None:
            return None
        box = provenance.bbox
        x0, x1 = sorted((float(box.l) / scale, float(box.r) / scale))
        y0, y1 = sorted((float(box.b) / scale, float(box.t) / scale))
        if x1 <= x0 or y1 <= y0:
            return None
        return BoundingBox(
            x0=x0,
            y0=y0,
            x1=x1,
            y1=y1,
            coordinate_space="points",
        )

    @staticmethod
    def _xlsx_range(provenance: ProvenanceItem) -> str:
        box = provenance.bbox
        start_column = max(1, math.floor(float(box.l)) + 1)
        end_column = max(start_column, math.ceil(float(box.r)))
        start_row = max(1, math.floor(min(float(box.t), float(box.b))) + 1)
        end_row = max(start_row, math.ceil(max(float(box.t), float(box.b))))
        start = f"{get_column_letter(start_column)}{start_row}"
        end = f"{get_column_letter(end_column)}{end_row}"
        return start if start == end else f"{start}:{end}"

    @staticmethod
    def _png_bytes(image: Image.Image) -> bytes:
        output = io.BytesIO()
        converted = image.convert("RGB")
        converted.save(output, format="PNG", optimize=False)
        return output.getvalue()

    @classmethod
    def _render_text_pages(
        cls,
        entries: Sequence[tuple[str, str]],
        *,
        title: str,
    ) -> tuple[dict[str, int], tuple[bytes, ...]]:
        """Render every DOCX element without truncation and return its start page."""
        title_font = cls._font(34, bold=True)
        body_font = cls._font(22)
        footer_font = cls._font(18)
        page_by_element: dict[str, int] = {}
        rendered_pages: list[bytes] = []
        page_number = 1

        def new_canvas(number: int) -> tuple[Image.Image, ImageDraw.ImageDraw, int]:
            canvas = Image.new("RGB", (1600, 2200), "white")
            draw = ImageDraw.Draw(canvas)
            draw.text((90, 70), title, fill="#12343B", font=title_font)
            draw.text(
                (90, 2_155),
                f"Normalized DOCX page {number}",
                fill="#52616B",
                font=footer_font,
            )
            return canvas, draw, 145

        canvas, draw, y = new_canvas(page_number)
        for element_key, content in entries:
            element_started = False
            for raw_line in content.splitlines() or [""]:
                for wrapped in textwrap.wrap(raw_line, width=105) or [""]:
                    if y + 31 > _TEXT_RENDER_BOTTOM:
                        rendered_pages.append(cls._png_bytes(canvas))
                        page_number += 1
                        canvas, draw, y = new_canvas(page_number)
                    if not element_started:
                        page_by_element[element_key] = page_number
                        element_started = True
                    draw.text((90, y), wrapped, fill="#1F2933", font=body_font)
                    y += 31
                y += 7
            y += 7
        rendered_pages.append(cls._png_bytes(canvas))
        return page_by_element, tuple(rendered_pages)

    @classmethod
    def _pptx_chart_renders(
        cls,
        source_path: Path,
    ) -> defaultdict[int, deque[tuple[str, bytes]]]:
        presentation = Presentation(str(source_path))
        result: defaultdict[int, deque[tuple[str, bytes]]] = defaultdict(deque)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_chart", False):
                    continue
                chart = shape.chart
                categories = [category.label for category in chart.plots[0].categories]
                series = [
                    (
                        str(item.name or f"Series {index}"),
                        [float(value or 0) for value in item.values],
                    )
                    for index, item in enumerate(chart.series, start=1)
                ]
                result[slide_number].append(
                    (
                        str(shape.name),
                        cls._render_chart(
                            title=f"Slide {slide_number} chart",
                            categories=categories,
                            series=series,
                        ),
                    )
                )
        return result

    @classmethod
    def _xlsx_chart_renders(
        cls,
        workbook: Any,
        values_workbook: Any,
    ) -> defaultdict[str, deque[tuple[str, bytes]]]:
        result: defaultdict[str, deque[tuple[str, bytes]]] = defaultdict(deque)
        for sheet in workbook.worksheets:
            for index, chart in enumerate(sheet._charts, start=1):
                categories: list[str] = []
                series: list[tuple[str, list[float]]] = []
                for series_index, item in enumerate(chart.ser, start=1):
                    value_ref = getattr(getattr(item.val, "numRef", None), "f", None)
                    category_ref = cls._chart_category_reference(item)
                    if not categories and category_ref:
                        categories = [
                            str(value)
                            for value in cls._workbook_reference_values(
                                values_workbook,
                                category_ref,
                            )
                        ]
                    values = [
                        cls._numeric_value(value)
                        for value in cls._workbook_reference_values(
                            values_workbook,
                            value_ref,
                        )
                    ]
                    title = getattr(item.tx, "v", None) or f"Series {series_index}"
                    series.append((str(title), values))
                result[sheet.title].append(
                    (
                        f"chart-{index}",
                        cls._render_chart(
                            title=cls._openpyxl_chart_title(chart) or f"Chart {index}",
                            categories=categories,
                            series=series,
                        ),
                    )
                )
        return result

    @staticmethod
    def _chart_category_reference(series: Any) -> str | None:
        category = series.cat
        for name in ("strRef", "numRef", "multiLvlStrRef"):
            candidate = getattr(category, name, None)
            reference = getattr(candidate, "f", None)
            if reference:
                return str(reference)
        return None

    @staticmethod
    def _workbook_reference_values(workbook: Any, reference: str | None) -> list[object]:
        if not reference:
            return []
        sheet_name, boundaries = range_to_tuple(reference)
        min_column, min_row, max_column, max_row = boundaries
        sheet = workbook[sheet_name]
        return [
            sheet.cell(row=row, column=column).value
            for row in range(min_row, max_row + 1)
            for column in range(min_column, max_column + 1)
        ]

    @staticmethod
    def _openpyxl_chart_title(chart: Any) -> str | None:
        try:
            paragraphs = chart.title.tx.rich.p
            return (
                " ".join(run.t for paragraph in paragraphs for run in paragraph.r if run.t).strip()
                or None
            )
        except (AttributeError, TypeError):
            return None

    @staticmethod
    def _workbook_manifest_and_formulas(
        workbook: Any,
        values_workbook: Any,
        filename: str,
    ) -> tuple[dict[str, object], list[ElementDraft]]:
        sheets: list[dict[str, object]] = []
        formulas: list[ElementDraft] = []
        for sheet in workbook.worksheets:
            formula_records: list[dict[str, object]] = []
            values_sheet = values_workbook[sheet.title]
            if sheet.sheet_state == "visible":
                for cell in sheet._cells.values():
                    if cell.data_type != "f":
                        continue
                    displayed = values_sheet[cell.coordinate].value
                    formula_records.append(
                        {
                            "cell": cell.coordinate,
                            "formula": cell.value,
                            "displayed_value": displayed,
                        }
                    )
                    formulas.append(
                        ElementDraft(
                            key=f"formula:{sheet.title}:{cell.coordinate}",
                            element_type="text",
                            original_content=(
                                f"Formula: {cell.value}\nDisplayed value: {displayed}"
                            ),
                            location=XlsxRangeLocation(
                                filename=filename,
                                sheet=sheet.title,
                                cell_range=cell.coordinate,
                            ),
                            extraction_method="openpyxl_formula_supplement_v1",
                        )
                    )
            sheets.append(
                {
                    "name": sheet.title,
                    "state": sheet.sheet_state,
                    "stored_nonempty_cells": len(sheet._cells),
                    "merged_ranges": sorted(str(item) for item in sheet.merged_cells.ranges),
                    "formula_records": formula_records,
                }
            )
        return {"filename": filename, "sheets": sheets}, formulas

    @classmethod
    def _render_chart(
        cls,
        *,
        title: str,
        categories: Sequence[str],
        series: Sequence[tuple[str, Sequence[float]]],
    ) -> bytes:
        width, height = 1400, 900
        canvas = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(canvas)
        title_font = cls._font(34, bold=True)
        label_font = cls._font(21)
        small_font = cls._font(18)
        draw.text((80, 50), title, fill="#12343B", font=title_font)
        left, top, right, bottom = 130, 165, 1320, 725
        draw.line((left, top, left, bottom), fill="#36454F", width=3)
        draw.line((left, bottom, right, bottom), fill="#36454F", width=3)
        maximum = max((value for _, values in series for value in values), default=1.0)
        maximum = max(maximum, 1.0)
        category_count = max(len(categories), 1)
        group_width = (right - left) / category_count
        bar_count = max(len(series), 1)
        bar_width = min(90.0, group_width * 0.72 / bar_count)
        for category_index in range(category_count):
            center = left + group_width * (category_index + 0.5)
            label = categories[category_index] if category_index < len(categories) else ""
            draw.text((center - 55, bottom + 20), str(label)[:14], fill="#1F2933", font=small_font)
            for series_index, (_name, values) in enumerate(series):
                value = values[category_index] if category_index < len(values) else 0.0
                x0 = center - (bar_count * bar_width) / 2 + series_index * bar_width
                x1 = x0 + bar_width * 0.82
                y0 = bottom - (bottom - top) * (value / maximum)
                draw.rectangle((x0, y0, x1, bottom), fill=_CHART_COLORS[series_index % 4])
                draw.text(
                    (x0 + 4, max(top, y0 - 27)),
                    f"{value:g}",
                    fill="#1F2933",
                    font=small_font,
                )
        legend_x = 110
        for index, (name, _values) in enumerate(series):
            draw.rectangle(
                (legend_x, 820, legend_x + 24, 844),
                fill=_CHART_COLORS[index % 4],
            )
            draw.text((legend_x + 34, 817), name[:24], fill="#1F2933", font=label_font)
            legend_x += 280
        return cls._png_bytes(canvas)

    @staticmethod
    def _font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
        candidates = (
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
        )
        for candidate in candidates:
            try:
                return ImageFont.truetype(candidate, size)
            except OSError:
                continue
        return ImageFont.load_default()

    @staticmethod
    def _numeric_value(value: object) -> float:
        if isinstance(value, int | float):
            return float(value)
        if isinstance(value, str):
            try:
                return float(value)
            except ValueError:
                return 0.0
        return 0.0
